from pptx import Presentation

from app.extractors.base_extractor import BaseExtractor


class PPTXExtractor(BaseExtractor):
    """
    Extract text from Microsoft PowerPoint (.pptx) presentations.
    """

    def extract(self, file_path: str) -> str:
        presentation = Presentation(file_path)

        slides_text = []

        for slide_number, slide in enumerate(presentation.slides, start=1):

            slide_content = [f"Slide {slide_number}"]

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        slide_content.append(text)

            slides_text.append("\n".join(slide_content))

        return "\n\n".join(slides_text)